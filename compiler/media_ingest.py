"""Ingest non-text raw sources: images, audio, and generic file attachments.

Images get an LLM-generated caption (vision-capable chat completion), audio
gets an LLM-generated transcript (speech-to-text completion), and both are
copied into wiki-app/static/media/ so they can be embedded in the compiled
page. Other file types are either text-extracted (PDF, DOCX, XLSX, PPTX via
their respective parsing libraries; CSV, TSV, JSON, XML, HTML, YAML, log, and
a ZIP's file listing via the standard library) or, when we have no parser for
them, registered as an opaque downloadable attachment — still copied to
static/media/ and linked, just without content extraction. That opaque
fallback is what lets the wiki accept practically any file format (legacy
office formats, archives beyond ZIP, video, ...) without needing a dedicated
parser for each one.

Every public function here returns plain dicts shaped like the chunk dicts
synthesizer.py already works with ({"chunk_index", "text", "source_type"}),
not RawChunk instances — that keeps this module import-free of synthesizer.py
and avoids a circular import.
"""

from __future__ import annotations

import csv
import hashlib
import json
from pathlib import Path

from models import STATIC_MEDIA_DIR
from text_chunking import split_text_into_chunks

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}

AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}

# File types we can pull text out of directly. Most need no LLM at all; DOCX/
# XLSX/PPTX need a parsing library (python-docx/openpyxl/python-pptx, all
# plain deps in requirements.txt) but still no LLM call for the extraction
# itself -- only the normal topic-extraction pass afterward, same as PDF.
TEXT_EXTRACTABLE_FILE_EXTENSIONS = {
    ".pdf",
    ".csv",
    ".json",
    ".tsv",
    ".xml",
    ".html",
    ".htm",
    ".yaml",
    ".yml",
    ".log",
    ".docx",
    ".xlsx",
    ".pptx",
    ".zip",
}

# Extensions with a real extractor (PDF/DOCX/XLSX/PPTX) whose failure to
# import/parse should be reported differently from formats with no
# extractor at all -- see the `note` logic in build_file_chunks().
_LIBRARY_EXTRACTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}

# File types we register as a downloadable attachment without parsing content
# -- archives (beyond ZIP's lightweight manifest), legacy/OpenDocument
# formats, and video, none of which have a text extractor here, but are
# still useful as linked/downloadable resources.
OPAQUE_FILE_EXTENSIONS = {
    ".rtf",
    ".odt",
    ".ods",
    ".odp",
    ".rar",
    ".7z",
    ".tar",
    ".gz",
    ".tgz",
    ".epub",
    ".mp4",
    ".mov",
    ".avi",
    ".mkv",
    ".m4v",
}

FILE_EXTENSIONS = TEXT_EXTRACTABLE_FILE_EXTENSIONS | OPAQUE_FILE_EXTENSIONS

IMAGE_CAPTION_SYSTEM_PROMPT = (
    "You are an assistant describing an image for a personal knowledge wiki. "
    "Describe what the image shows in 2-4 sentences: subject, any visible "
    "text, diagrams, charts, or notable details. Be factual and concise. If "
    "the image is a screenshot, transcribe any clearly readable text."
)

FILE_CONTENT_MAX_CHARS = 6000


def copy_bytes_to_static(data: bytes, filename: str, static_dir: Path | None = None) -> Path:
    """Write bytes into wiki-app/static/media/, deduped by content hash."""
    out_dir = static_dir or STATIC_MEDIA_DIR
    out_dir.mkdir(parents=True, exist_ok=True)

    digest = hashlib.sha256(data).hexdigest()[:10]
    stem = Path(filename).stem or "file"
    suffix = Path(filename).suffix.lower()
    dest = out_dir / f"{stem}-{digest}{suffix}"
    if not dest.exists():
        dest.write_bytes(data)
    return dest


def copy_media_to_static(source: Path, static_dir: Path | None = None) -> Path:
    """Copy a raw media file on disk into wiki-app/static/media/, deduped by hash."""
    return copy_bytes_to_static(source.read_bytes(), source.name, static_dir)


def docs_relative_media_link(dest_path: Path, static_dir: Path | None = None) -> str:
    """Build a markdown-relative link from a wiki-app/docs/*.md page to a static asset.

    wiki-app/docs/ and wiki-app/static/ are sibling directories, so a docs page
    reaches a file at wiki-app/static/media/x.png via "../static/media/x.png".
    Docusaurus's markdown image/link transform resolves relative paths like
    this at build time and rewrites them correctly for the configured baseUrl
    (unlike a leading "/" absolute path, which is NOT baseUrl-prefixed and so
    breaks between local dev and a GitHub Pages subpath deploy).
    """
    out_dir = static_dir or STATIC_MEDIA_DIR
    wiki_app_dir = out_dir.parent.parent  # .../wiki-app/static/media -> .../wiki-app
    rel = dest_path.resolve().relative_to(wiki_app_dir.resolve())
    return f"../{rel.as_posix()}"


def describe_image(path: Path, llm) -> str:
    """Ask the LLM (vision-capable) to caption an image."""
    return llm.describe_image(path, IMAGE_CAPTION_SYSTEM_PROMPT).strip()


def build_image_chunk(path: Path, rel_source: str, llm, static_dir: Path | None = None) -> dict:
    """Build a single chunk dict for an image: caption text + embedded markdown image."""
    dest = copy_media_to_static(path, static_dir)
    link = docs_relative_media_link(dest, static_dir)
    caption = describe_image(path, llm)

    text = (
        f"Image file: `{rel_source}`\n\n"
        f"Description: {caption}\n\n"
        f"![{path.stem}]({link})"
    )
    return {"chunk_index": 0, "text": text, "source_type": "image", "media_link": link}


def transcribe_audio(path: Path, llm) -> str:
    """Ask the LLM (speech-to-text) to transcribe an audio file."""
    return llm.transcribe_audio(path).strip()


def build_audio_chunk(path: Path, rel_source: str, llm, static_dir: Path | None = None) -> dict:
    """Build a single chunk dict for an audio file: transcript text + player link.

    Unlike images (which require an LLM -- see synthesizer.py's require_llm
    call), transcription degrades gracefully: no LLM configured, or the
    transcription call itself failing, falls back to a metadata + download
    chunk instead of crashing the compile, mirroring the PDF text-extraction
    fallback below.
    """
    transcript = ""
    if llm is not None and getattr(llm, "available", False):
        try:
            transcript = transcribe_audio(path, llm)
        except Exception:
            transcript = ""

    dest = copy_media_to_static(path, static_dir)
    link = docs_relative_media_link(dest, static_dir)

    if transcript:
        text = (
            f"Audio file: `{rel_source}`\n\n"
            f"Transcript: {transcript}\n\n"
            f"[Listen to {path.name}]({link})"
        )
    else:
        size_kb = path.stat().st_size / 1024
        text = (
            f"Audio file: `{rel_source}` ({path.suffix.lstrip('.').upper()}, {size_kb:.1f} KB) "
            f"— transcription unavailable\n\n"
            f"[Listen to {path.name}]({link})"
        )
    return {"chunk_index": 0, "text": text, "source_type": "audio", "media_link": link}


def _extract_pdf_text(path: Path) -> str | None:
    """Best-effort PDF text extraction. Returns None if pypdf (or one of its
    own dependencies) isn't usable in this environment — deliberately broad
    because a broken transitive dependency (seen in practice: pypdf's
    optional `cryptography` extra failing to import) doesn't raise a plain
    ImportError, and a PDF attachment we can't parse should degrade to an
    opaque attachment rather than crash the whole compile."""
    try:
        from pypdf import PdfReader
    except BaseException:
        # Deliberately catches more than ImportError/Exception: a broken
        # transitive dependency of pypdf (observed in practice — a system
        # `cryptography` install missing its native backend) surfaces as a
        # pyo3_runtime.PanicException, which does not subclass Exception.
        # This is best-effort text extraction; any failure to even import
        # the library should fall back to treating the PDF as an opaque
        # attachment, never crash the compile.
        return None

    try:
        reader = PdfReader(str(path))
        pages_text = [page.extract_text() or "" for page in reader.pages]
    except Exception:
        return None

    text = "\n\n".join(t.strip() for t in pages_text if t.strip())
    return text.strip() or None


def _extract_docx_text(path: Path) -> str | None:
    """Best-effort DOCX text extraction via python-docx: paragraphs in
    document order plus table cells. Returns None (same contract as
    _extract_pdf_text above) if the library isn't usable or the file can't
    be parsed, so an unparseable DOCX degrades to an opaque attachment
    rather than crashing the compile."""
    try:
        import docx
    except BaseException:
        return None

    try:
        document = docx.Document(str(path))
        parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
    except Exception:
        return None

    text = "\n".join(parts).strip()
    return text or None


def _extract_xlsx_text(path: Path, *, max_rows_per_sheet: int = 100) -> str | None:
    """Best-effort XLSX text extraction via openpyxl: every sheet's cell
    values as comma-joined rows, capped per sheet so a huge spreadsheet
    doesn't blow up chunk size. Returns None on any import/parse failure,
    same opaque-fallback contract as _extract_pdf_text/_extract_docx_text."""
    try:
        import openpyxl
    except BaseException:
        return None

    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sections = []
        for sheet in workbook.worksheets:
            lines = [f"## Sheet: {sheet.title}"]
            row_count = 0
            for row_count, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_count > max_rows_per_sheet:
                    lines.append(f"… (more rows not shown, sheet has {sheet.max_row} total)")
                    break
                cells = ["" if cell is None else str(cell).strip() for cell in row]
                if any(cells):
                    lines.append(", ".join(cells))
            if row_count:
                sections.append("\n".join(lines))
        workbook.close()
    except Exception:
        return None

    text = "\n\n".join(sections).strip()
    return text or None


def _extract_pptx_text(path: Path) -> str | None:
    """Best-effort PPTX text extraction via python-pptx: each slide's shape
    text, table cells, and speaker notes, in slide order. Returns None on
    any import/parse failure, same opaque-fallback contract as the other
    _extract_*_text helpers above."""
    try:
        from pptx import Presentation
    except BaseException:
        return None

    try:
        presentation = Presentation(str(path))
        slides_text = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines = [f"## Slide {index}"]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            lines.append(line)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            lines.append(" | ".join(cells))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"Notes: {notes}")
            slides_text.append("\n".join(lines))
    except Exception:
        return None

    text = "\n\n".join(slides_text).strip()
    return text or None


def _extract_zip_manifest(path: Path, *, max_entries: int = 200) -> str | None:
    """List a ZIP archive's contents (name + size) as a lightweight manifest.

    Deliberately not full recursive extraction -- ingesting an archive's
    contents as first-class raw sources would need real sandboxing against
    zip bombs and path traversal, plus its own place in the incremental
    state/dedup model, which is a bigger feature than "make an archive
    searchable." A manifest is enough to know what's inside without
    downloading it, and is still a real improvement over a bare download
    link with zero information about what's in the archive.
    """
    import zipfile

    try:
        with zipfile.ZipFile(path) as archive:
            infos = [info for info in archive.infolist() if not info.is_dir()]
    except (zipfile.BadZipFile, OSError):
        return None

    lines = [f"{info.filename} ({info.file_size / 1024:.1f} KB)" for info in infos[:max_entries]]
    remaining = len(infos) - max_entries
    if remaining > 0:
        lines.append(f"… ({remaining} more {'entry' if remaining == 1 else 'entries'} not shown)")

    text = "\n".join(lines).strip()
    return text or None


def _extract_delimited_text(path: Path, *, delimiter: str = ",", max_rows: int = 50) -> str:
    with path.open(newline="", encoding="utf-8", errors="replace") as handle:
        rows = list(csv.reader(handle, delimiter=delimiter))
    preview = rows[:max_rows]
    lines = [", ".join(cell.strip() for cell in row) for row in preview]
    note = ""
    if len(rows) > max_rows:
        note = f"\n\n… ({len(rows) - max_rows} more row(s) not shown)"
    return "\n".join(lines) + note


def _extract_csv_text(path: Path, *, max_rows: int = 50) -> str:
    return _extract_delimited_text(path, delimiter=",", max_rows=max_rows)


def _extract_tsv_text(path: Path, *, max_rows: int = 50) -> str:
    return _extract_delimited_text(path, delimiter="\t", max_rows=max_rows)


def _extract_json_text(path: Path) -> str:
    data = json.loads(path.read_text(encoding="utf-8"))
    return json.dumps(data, indent=2, ensure_ascii=False)[:FILE_CONTENT_MAX_CHARS]


def _extract_plain_text_file(path: Path) -> str:
    """Best-effort raw-text extraction for markup/config formats (XML, HTML,
    YAML, log files) that need no dedicated parser -- the file's own content
    already is the text worth chunking."""
    return path.read_text(encoding="utf-8", errors="replace")[:FILE_CONTENT_MAX_CHARS]


def _opaque_file_chunk(path: Path, rel_source: str, static_dir: Path | None, note: str = "") -> dict:
    dest = copy_media_to_static(path, static_dir)
    link = docs_relative_media_link(dest, static_dir)
    size_kb = path.stat().st_size / 1024
    text = (
        f"Attached file: `{rel_source}` ({path.suffix.lstrip('.').upper()}, {size_kb:.1f} KB)"
        f"{f' — {note}' if note else ''}\n\n"
        f"[Download {path.name}]({link})"
    )
    return {"chunk_index": 0, "text": text, "source_type": "file", "media_link": link}


def build_file_chunks(path: Path, rel_source: str, static_dir: Path | None = None) -> list[dict]:
    """Build chunk dict(s) for a generic file attachment.

    Everything in TEXT_EXTRACTABLE_FILE_EXTENSIONS (PDF, DOCX, XLSX, PPTX,
    CSV, TSV, JSON, XML, HTML, YAML, log, ZIP) gets its text (or, for ZIP, its
    file listing) extracted and chunked like any raw text source; everything
    else in FILE_EXTENSIONS becomes a single metadata + download-link chunk
    with no content extraction. A PDF/DOCX/XLSX/PPTX whose parsing library
    isn't installed, or whose file fails to parse, degrades to that same
    opaque fallback rather than crashing the compile.
    """
    suffix = path.suffix.lower()

    extracted: str | None = None
    if suffix == ".pdf":
        extracted = _extract_pdf_text(path)
    elif suffix == ".csv":
        extracted = _extract_csv_text(path)
    elif suffix == ".tsv":
        extracted = _extract_tsv_text(path)
    elif suffix == ".json":
        extracted = _extract_json_text(path)
    elif suffix in {".xml", ".html", ".htm", ".yaml", ".yml", ".log"}:
        extracted = _extract_plain_text_file(path)
    elif suffix == ".docx":
        extracted = _extract_docx_text(path)
    elif suffix == ".xlsx":
        extracted = _extract_xlsx_text(path)
    elif suffix == ".pptx":
        extracted = _extract_pptx_text(path)
    elif suffix == ".zip":
        extracted = _extract_zip_manifest(path)

    if extracted:
        dest = copy_media_to_static(path, static_dir)
        link = docs_relative_media_link(dest, static_dir)
        label = "contents" if suffix == ".zip" else "download"
        header = f"Attached file: `{rel_source}` ([{label}]({link}))\n\n"
        pieces = split_text_into_chunks(extracted, max_chars=FILE_CONTENT_MAX_CHARS)
        if not pieces:
            pieces = [""]
        return [
            {
                "chunk_index": index,
                "text": header + piece if index == 0 else piece,
                "source_type": "file",
                "media_link": link,
            }
            for index, piece in enumerate(pieces)
        ]

    if suffix in _LIBRARY_EXTRACTED_EXTENSIONS:
        note = "text extraction unavailable"
    elif suffix == ".zip":
        note = "could not read archive contents"
    else:
        note = "content not parsed"
    return [_opaque_file_chunk(path, rel_source, static_dir, note=note)]
